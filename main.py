from fastapi import FastAPI, Query, Request, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import RedirectResponse, JSONResponse
from google.oauth2.credentials import Credentials
from google.analytics.admin import AnalyticsAdminServiceClient
import requests
import os
import json
import base64
from ga4audit import run_ga4_audit_with_creds

app = FastAPI()

# ── Environment Variables ──────────────────────────────────────────────────
GOOGLE_CLIENT_ID     = os.getenv("GOOGLE_CLIENT_ID")
GOOGLE_CLIENT_SECRET = os.getenv("GOOGLE_CLIENT_SECRET")
REDIRECT_URI         = os.getenv("REDIRECT_URI", "https://ga4-audit-backend.onrender.com/auth/callback")
FRONTEND_URL         = os.getenv("FRONTEND_URL", "https://ga4-audit-frontend.vercel.app")

GOOGLE_AUTH_URL  = "https://accounts.google.com/o/oauth2/v2/auth"
GOOGLE_TOKEN_URL = "https://oauth2.googleapis.com/token"

SCOPES = [
    "https://www.googleapis.com/auth/analytics.readonly",
    "https://www.googleapis.com/auth/analytics.edit",
    "openid",
    "email",
    "profile",
]

# ── CORS ───────────────────────────────────────────────────────────────────
# Must be registered BEFORE any route handlers
app.add_middleware(
    CORSMiddleware,
    allow_origins=[
        FRONTEND_URL,
        "https://ga4-audit-frontend.vercel.app",
        "http://localhost:5173",
        "http://localhost:5174",
        "http://localhost:5175",
    ],
    allow_credentials=True,
    allow_methods=["GET", "POST", "PUT", "DELETE", "OPTIONS", "PATCH"],
    allow_headers=["*"],
    expose_headers=["Content-Disposition", "Content-Type", "Content-Length"],
    max_age=3600,
)

# ── Explicit OPTIONS handler for all routes (handles Render cold-start) ────
@app.options("/{rest_of_path:path}")
async def preflight_handler(request: Request, rest_of_path: str):
    return JSONResponse(
        content={},
        headers={
            "Access-Control-Allow-Origin":  request.headers.get("origin", "*"),
            "Access-Control-Allow-Methods": "GET, POST, PUT, DELETE, OPTIONS, PATCH",
            "Access-Control-Allow-Headers": "*",
            "Access-Control-Allow-Credentials": "true",
            "Access-Control-Max-Age": "3600",
        },
    )

# ── Health / keep-alive endpoints ─────────────────────────────────────────
@app.get("/")
def root():
    return {"status": "ok", "service": "GA4 Audit Backend"}

@app.get("/health")
def health():
    return {"status": "healthy"}

# ── Helper: decode token from request header ───────────────────────────────
def get_user_credentials(request: Request) -> Credentials:
    """
    Frontend sends: Authorization: Bearer <base64-encoded-token-json>
    """
    auth_header = request.headers.get("Authorization", "")
    if not auth_header.startswith("Bearer "):
        raise HTTPException(status_code=401, detail="Not authenticated. Please log in.")
    try:
        token_json = base64.b64decode(auth_header[7:]).decode("utf-8")
        token_data = json.loads(token_json)
    except Exception:
        raise HTTPException(status_code=401, detail="Invalid token format.")

    return Credentials(
        token=token_data["access_token"],
        refresh_token=token_data.get("refresh_token"),
        token_uri=GOOGLE_TOKEN_URL,
        client_id=GOOGLE_CLIENT_ID,
        client_secret=GOOGLE_CLIENT_SECRET,
        scopes=SCOPES,
    )

# ── Routes ─────────────────────────────────────────────────────────────────

@app.get("/")
def read_root():
    return {"message": "GA4 Audit API is running 🚀"}


@app.get("/auth/google")
def auth_google():
    """Step 1 — redirect user to Google consent screen."""
    scope_str = " ".join(SCOPES)
    params = (
        f"?client_id={GOOGLE_CLIENT_ID}"
        f"&redirect_uri={REDIRECT_URI}"
        f"&response_type=code"
        f"&scope={requests.utils.quote(scope_str)}"
        f"&access_type=offline"
        f"&prompt=consent"
    )
    return RedirectResponse(url=GOOGLE_AUTH_URL + params)


@app.get("/auth/callback")
def auth_callback(code: str = Query(...)):
    """
    Step 2 — exchange code for tokens, then redirect to frontend
    with the token data base64-encoded in the URL fragment (#).
    Fragments never leave the browser so the token stays client-side only.
    """
    token_response = requests.post(GOOGLE_TOKEN_URL, data={
        "code": code,
        "client_id": GOOGLE_CLIENT_ID,
        "client_secret": GOOGLE_CLIENT_SECRET,
        "redirect_uri": REDIRECT_URI,
        "grant_type": "authorization_code",
    })

    if token_response.status_code != 200:
        return RedirectResponse(url=f"{FRONTEND_URL}?auth=error&reason=token_exchange_failed")

    token_data = token_response.json()

    # Fetch user info
    userinfo_resp = requests.get(
        "https://www.googleapis.com/oauth2/v3/userinfo",
        headers={"Authorization": f"Bearer {token_data['access_token']}"}
    )
    user_info = userinfo_resp.json() if userinfo_resp.status_code == 200 else {}

    # Bundle token + user info, base64 encode, pass via URL fragment
    payload = json.dumps({
        "token_data": token_data,
        "user_info": user_info,
    })
    encoded = base64.urlsafe_b64encode(payload.encode()).decode()

    # Use fragment (#) — never sent to server, stays in browser only
    return RedirectResponse(url=f"{FRONTEND_URL}#auth={encoded}")


@app.get("/list-properties")
def list_properties(request: Request):
    """Returns all GA4 properties the logged-in user has access to."""
    creds = get_user_credentials(request)
    try:
        admin_client = AnalyticsAdminServiceClient(credentials=creds)
        accounts = admin_client.list_account_summaries()
        properties = []
        for account in accounts:
            for prop in account.property_summaries:
                properties.append({
                    "property_id": prop.property.replace("properties/", ""),
                    "display_name": prop.display_name,
                    "account_name": account.display_name,
                })
        return {"success": True, "properties": properties}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.get("/run-audit")
def run_audit(
    request: Request,
    property_id: str = Query(...),
    start_date: str = Query("30daysAgo"),
    end_date: str = Query("today"),
):
    """Runs the full GA4 audit using the logged-in user's credentials."""
    creds = get_user_credentials(request)
    try:
        results = run_ga4_audit_with_creds(
            creds=creds,
            property_numeric_id=property_id,
            start_date=start_date,
            end_date=end_date,
        )
        return {"success": True, "data": results}
    except Exception as e:
        return {"success": False, "error": str(e)}


@app.post("/sdr-report")
def sdr_report(
    request: Request,
    property_id: str = Query(...),
    start_date: str = Query("30daysAgo"),
    end_date: str = Query("today"),
    body: dict = None,
):
    """
    Fetches live GA4 data for a list of SDR events.
    POST body: { "events": ["event_name_1", "event_name_2", ...], "params": ["param1", "param2"] }
    Returns event counts and top parameter values for each requested event.
    """
    from google.analytics.data_v1beta import BetaAnalyticsDataClient
    from google.analytics.data_v1beta.types import (
        RunReportRequest, Dimension, Metric, FilterExpression,
        Filter, FilterExpressionList
    )

    creds = get_user_credentials(request)

    if not body:
        return {"success": False, "error": "Request body required with 'events' list."}

    event_names = body.get("events", [])
    param_names = body.get("params", [])

    if not event_names:
        return {"success": False, "error": "No events provided."}

    try:
        data_client = BetaAnalyticsDataClient(credentials=creds)
        prop = f"properties/{property_id}"
        report_rows = []

        # ── Step 1: Event counts per event name ────────────────────────────
        event_count_req = RunReportRequest(
            property=prop,
            dimensions=[Dimension(name="eventName")],
            metrics=[Metric(name="eventCount"), Metric(name="totalUsers")],
            date_ranges=[{"start_date": start_date, "end_date": end_date}],
            dimension_filter=FilterExpression(
                or_group=FilterExpressionList(
                    expressions=[
                        FilterExpression(filter=Filter(
                            field_name="eventName",
                            string_filter=Filter.StringFilter(
                                match_type=Filter.StringFilter.MatchType.EXACT,
                                value=evt,
                            )
                        ))
                        for evt in event_names
                    ]
                )
            ),
            limit=1000,
        )
        event_counts = {}
        event_users  = {}
        resp = data_client.run_report(request=event_count_req)
        for row in resp.rows:
            evt = row.dimension_values[0].value
            event_counts[evt] = int(row.metric_values[0].value)
            event_users[evt]  = int(row.metric_values[1].value)

        # ── Step 2: Per-event, per-param top values ────────────────────────
        param_data = {}  # { "event_name": { "param_name": [top values] } }

        for evt in event_names:
            param_data[evt] = {}
            for param in param_names:
                if not param:
                    continue
                # Use eventName + customEvent:param or just param dimension
                try:
                    param_req = RunReportRequest(
                        property=prop,
                        dimensions=[
                            Dimension(name="eventName"),
                            Dimension(name=f"customEvent:{param}"),
                        ],
                        metrics=[Metric(name="eventCount")],
                        date_ranges=[{"start_date": start_date, "end_date": end_date}],
                        dimension_filter=FilterExpression(
                            filter=Filter(
                                field_name="eventName",
                                string_filter=Filter.StringFilter(
                                    match_type=Filter.StringFilter.MatchType.EXACT,
                                    value=evt,
                                )
                            )
                        ),
                        limit=10,
                        order_bys=[{"metric": {"metric_name": "eventCount"}, "desc": True}],
                    )
                    param_resp = data_client.run_report(request=param_req)
                    top_values = []
                    for row in param_resp.rows:
                        val = row.dimension_values[1].value
                        cnt = int(row.metric_values[0].value)
                        if val and val not in ["(not set)", ""]:
                            top_values.append({"value": val, "count": cnt})
                    param_data[evt][param] = top_values
                except Exception:
                    param_data[evt][param] = []

        # ── Build response ─────────────────────────────────────────────────
        for evt in event_names:
            report_rows.append({
                "eventName":   evt,
                "eventCount":  event_counts.get(evt, 0),
                "totalUsers":  event_users.get(evt, 0),
                "inGA4":       evt in event_counts,
                "paramData":   param_data.get(evt, {}),
            })

        return {"success": True, "report": report_rows}

    except Exception as e:
        return {"success": False, "error": str(e)}

@app.post("/explore")
def explore(
    request: Request,
    property_id: str = Query(...),
    start_date: str = Query("30daysAgo"),
    end_date: str = Query("today"),
    body: dict = None,
):
    """
    Flexible GA4 Data API explorer.
    POST body: {
      "dimensions": ["eventName", "sessionDefaultChannelGroup", ...],
      "metrics": ["eventCount", "totalUsers", "sessions", ...],
      "limit": 100,
      "order_by_metric": "eventCount"
    }
    Returns rows with dimension values and metric values.
    """
    from google.analytics.data_v1beta import BetaAnalyticsDataClient
    from google.analytics.data_v1beta.types import (
        RunReportRequest, Dimension, Metric, OrderBy
    )

    creds = get_user_credentials(request)

    if not body:
        return {"success": False, "error": "Request body required."}

    dimensions   = body.get("dimensions", [])
    metrics      = body.get("metrics", [])
    limit        = min(int(body.get("limit", 100)), 5000)
    order_metric = body.get("order_by_metric", metrics[0] if metrics else None)

    if not dimensions and not metrics:
        return {"success": False, "error": "At least one dimension or metric required."}

    try:
        data_client = BetaAnalyticsDataClient(credentials=creds)

        order_bys = []
        if order_metric and order_metric in metrics:
            order_bys = [OrderBy(
                metric=OrderBy.MetricOrderBy(metric_name=order_metric),
                desc=True,
            )]

        # ── Build dimension filter from body ───────────────────────────────
        filters_list = body.get("filters", [])  # [{dimension, matchType, value}]
        dimension_filter = None

        if filters_list:
            from google.analytics.data_v1beta.types import FilterExpression, FilterExpressionList, Filter

            def build_filter(f):
                match_map = {
                    "EXACT":       Filter.StringFilter.MatchType.EXACT,
                    "BEGINS_WITH": Filter.StringFilter.MatchType.BEGINS_WITH,
                    "ENDS_WITH":   Filter.StringFilter.MatchType.ENDS_WITH,
                    "CONTAINS":    Filter.StringFilter.MatchType.CONTAINS,
                    "REGEXP":      Filter.StringFilter.MatchType.FULL_REGEXP,
                }
                return FilterExpression(filter=Filter(
                    field_name=f["dimension"],
                    string_filter=Filter.StringFilter(
                        match_type=match_map.get(f.get("matchType","CONTAINS"), Filter.StringFilter.MatchType.CONTAINS),
                        value=f["value"],
                        case_sensitive=False,
                    )
                ))

            if len(filters_list) == 1:
                dimension_filter = build_filter(filters_list[0])
            else:
                dimension_filter = FilterExpression(
                    and_group=FilterExpressionList(
                        expressions=[build_filter(f) for f in filters_list]
                    )
                )

        req = RunReportRequest(
            property=f"properties/{property_id}",
            dimensions=[Dimension(name=d) for d in dimensions],
            metrics=[Metric(name=m) for m in metrics],
            date_ranges=[{"start_date": start_date, "end_date": end_date}],
            limit=limit,
            order_bys=order_bys,
            dimension_filter=dimension_filter,
        )

        resp = data_client.run_report(request=req)

        rows = []
        for row in resp.rows:
            r = {}
            for i, d in enumerate(dimensions):
                r[d] = row.dimension_values[i].value
            for i, m in enumerate(metrics):
                r[m] = row.metric_values[i].value
            rows.append(r)

        return {
            "success": True,
            "rows": rows,
            "row_count": len(rows),
            "dimensions": dimensions,
            "metrics": metrics,
        }

    except Exception as e:
        return {"success": False, "error": str(e)}

@app.get("/custom-dimensions")
def get_custom_dimensions(
    request: Request,
    property_id: str = Query(...),
):
    """Returns all custom dimensions for a property to populate the Explorer dimension list."""
    from google.analytics.admin import AnalyticsAdminServiceClient
    creds = get_user_credentials(request)
    try:
        admin_client = AnalyticsAdminServiceClient(credentials=creds)
        dims = list(admin_client.list_custom_dimensions(parent=f"properties/{property_id}"))
        result = []
        for d in dims:
            scope_int = int(d.scope)
            prefix = {1: "customEvent", 2: "customUser", 3: "customItem"}.get(scope_int, "customEvent")
            scope_label = {1: "Event", 2: "User", 3: "Item"}.get(scope_int, "Event")
            result.append({
                "id":           f"{prefix}:{d.parameter_name}",
                "label":        f"{d.display_name} ({scope_label})",
                "paramName":    d.parameter_name,
                "displayName":  d.display_name,
                "scope":        scope_label,
                "group":        f"Custom — {scope_label} Scoped",
            })
        return {"success": True, "dimensions": result}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/export/xlsx")
def export_xlsx(
    request: Request,
    body: dict = None,
):
    """
    Generates an audit report as a multi-sheet .xlsx file.
    POST body: { "data": <audit_data_object>, "property_name": "...", "date_range": "..." }
    """
    import io
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
    from openpyxl.utils import get_column_letter
    from fastapi.responses import StreamingResponse

    get_user_credentials(request)  # auth check

    if not body:
        return {"success": False, "error": "Request body required."}

    audit_data   = body.get("data", {})
    prop_name    = body.get("property_name", "GA4 Property")
    date_range   = body.get("date_range", "")

    # ── Brainlabs colours ──────────────────────────────────────────────────
    BL_BLACK  = "0A0A0A"
    BL_YELLOW = "FFD426"
    BL_WHITE  = "FFFFFF"
    BL_DGREY  = "1A1A1A"
    BL_MGREY  = "3D3D3D"
    BL_LGREY  = "9A9A9A"
    BL_GREEN  = "00C896"
    BL_RED    = "FF4444"

    wb = Workbook()
    wb.remove(wb.active)  # remove default sheet

    def make_header_style(ws, row, cols, bg=BL_BLACK, fg=BL_YELLOW, bold=True, size=11):
        for col in range(1, cols+1):
            c = ws.cell(row=row, column=col)
            c.font = Font(name="Arial", bold=bold, color=fg, size=size)
            c.fill = PatternFill("solid", fgColor=bg)
            c.alignment = Alignment(horizontal="left", vertical="center", wrap_text=False)

    def make_row_style(ws, row, cols, even=False):
        bg = "F7F6F2" if even else BL_WHITE
        for col in range(1, cols+1):
            c = ws.cell(row=row, column=col)
            c.font = Font(name="Arial", size=10, color=BL_BLACK)
            c.fill = PatternFill("solid", fgColor=bg)
            c.alignment = Alignment(horizontal="left", vertical="center", wrap_text=True)

    def set_col_widths(ws, widths):
        for i, w in enumerate(widths, 1):
            ws.column_dimensions[get_column_letter(i)].width = w

    def add_title_row(ws, title, subtitle=""):
        ws.row_dimensions[1].height = 32
        ws["A1"] = title
        ws["A1"].font = Font(name="Arial", bold=True, size=14, color=BL_YELLOW)
        ws["A1"].fill = PatternFill("solid", fgColor=BL_BLACK)
        ws["A1"].alignment = Alignment(horizontal="left", vertical="center")
        if subtitle:
            ws["B1"] = subtitle
            ws["B1"].font = Font(name="Arial", size=10, color=BL_LGREY)
            ws["B1"].fill = PatternFill("solid", fgColor=BL_BLACK)
            ws["B1"].alignment = Alignment(horizontal="left", vertical="center")

    # ── Sheet builder helper ───────────────────────────────────────────────
    def add_check_result_sheet(name, entries):
        ws = wb.create_sheet(title=name[:31])
        add_title_row(ws, name, f"{prop_name} · {date_range}")
        ws.merge_cells("A1:B1")
        ws.row_dimensions[2].height = 22
        ws["A2"] = "Check"; ws["B2"] = "Result"
        make_header_style(ws, 2, 2)
        for i, e in enumerate(entries):
            r = i + 3
            result_val = e.get("Result", "")
            if isinstance(result_val, (list, dict)):
                result_val = str(result_val)
            ws.cell(row=r, column=1).value = e.get("Check", "")
            ws.cell(row=r, column=2).value = result_val
            make_row_style(ws, r, 2, even=(i % 2 == 0))
            # Colour-code results
            c = ws.cell(row=r, column=2)
            if isinstance(result_val, str):
                if "✅" in result_val:
                    c.font = Font(name="Arial", size=10, color=BL_GREEN)
                elif "❌" in result_val:
                    c.font = Font(name="Arial", size=10, color=BL_RED)
        set_col_widths(ws, [42, 55])

    def add_table_sheet(name, headers, rows_data):
        ws = wb.create_sheet(title=name[:31])
        add_title_row(ws, name, f"{prop_name} · {date_range}")
        ws.merge_cells(f"A1:{get_column_letter(len(headers))}1")
        hr = 2
        ws.row_dimensions[hr].height = 22
        for ci, h in enumerate(headers, 1):
            ws.cell(row=hr, column=ci).value = h
        make_header_style(ws, hr, len(headers))
        for i, row in enumerate(rows_data):
            r = i + 3
            for ci, val in enumerate(row, 1):
                ws.cell(row=r, column=ci).value = val
            make_row_style(ws, r, len(headers), even=(i % 2 == 0))
        col_w = max(18, 60 // max(len(headers), 1))
        set_col_widths(ws, [col_w] * len(headers))

    # ── Build sheets ───────────────────────────────────────────────────────
    sections = [
        ("Property Details",     audit_data.get("Property Details", [])),
        ("Streams Config",       audit_data.get("Streams Configuration", [])),
        ("GA4 Property Limits",  audit_data.get("GA4 Property Limits", [])),
        ("PII Check",            audit_data.get("PII Check", [])),
        ("Transactions",         audit_data.get("Transactions", [])),
        ("Landing Page Analysis",audit_data.get("Landing Page Analysis", [])),
        ("Channel Analysis",     audit_data.get("Channel Grouping Analysis", [])),
    ]
    for name, entries in sections:
        if entries:
            add_check_result_sheet(name, entries)

    # Events
    events = audit_data.get("GA4 Events", [])
    if events:
        add_table_sheet("Event Inventory", ["Event Name", "Event Count"],
                        [(e.get("Check",""), e.get("Result","")) for e in events])

    # Custom Dimensions — all scopes
    for scope in ["Event Scoped", "User Scoped", "Item Scoped"]:
        dims = audit_data.get(f"Custom Dimensions - {scope}", [])
        if dims:
            add_table_sheet(f"Custom Dims ({scope[:4]})", ["Display Name","Parameter Name","Scope","Ads Personalisation"],
                [(e.get("Check",""), e.get("Result",{}).get("Parameter Name",""), e.get("Result",{}).get("Scope",""), e.get("Result",{}).get("Ads Personalization Excluded","")) for e in dims])

    # Key Events
    key_events = audit_data.get("Key Event Details", [])
    if key_events:
        add_table_sheet("Key Events", ["Event Name","Create Time","Counting Method"],
                [(e.get("Check",""), e.get("Result",{}).get("Create Time",""), e.get("Result",{}).get("Counting Method","")) for e in key_events])

    # Duplicate Transactions
    dup_txns = audit_data.get("Duplicate Transactions", [])
    if dup_txns:
        add_table_sheet("Duplicate Transactions", ["Transaction ID","Count"],
                        [(e.get("transactionId",""), e.get("count","")) for e in dup_txns])

    # Landing Page Data
    lp_data = audit_data.get("Landing Page Data", [])
    if lp_data:
        add_table_sheet("Landing Page Data", ["Landing Page","Sessions"],
                        [(e.get("Landing Page",""), e.get("Sessions","")) for e in lp_data[:200]])

    # Channel Data
    ch_data = audit_data.get("Channel Grouping Data", [])
    if ch_data:
        add_table_sheet("Channel Grouping", ["Channel Group","Sessions"],
                        [(e.get("Channel Group",""), e.get("Sessions","")) for e in ch_data])

    # Unassigned source/medium
    ua_data = audit_data.get("Unassigned Source/Medium Data", [])
    if ua_data:
        add_table_sheet("Unassigned Traffic", ["Channel","Source","Medium","Sessions"],
                        [(e.get("Channel Group",""),e.get("Source",""),e.get("Medium",""),e.get("Sessions","")) for e in ua_data])

    # ── Stream to response ─────────────────────────────────────────────────
    buf = io.BytesIO()
    wb.save(buf)
    buf.seek(0)
    filename = f"GA4_Audit_{prop_name.replace(' ','_')}.xlsx"
    return StreamingResponse(
        buf,
        media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        headers={"Content-Disposition": f'attachment; filename="{filename}"'},
    )

@app.post("/export/pptx")
def export_pptx(
    request: Request,
    body: dict = None,
):
    """
    Generates a Brainlabs-branded GA4 audit .pptx using pure python-pptx.
    POST body: { "data": <audit_data>, "property_name": "...", "date_range": "..." }
    """
    import io
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt
    from fastapi.responses import StreamingResponse

    get_user_credentials(request)

    if not body:
        return {"success": False, "error": "Request body required."}

    audit_data = body.get("data", {})
    prop_name  = body.get("property_name", "GA4 Property")
    date_range = body.get("date_range", "30daysAgo – today")

    # ── Brainlabs colours ──────────────────────────────────────────────────
    C_BLACK  = RGBColor(0x0A, 0x0A, 0x0A)
    C_YELLOW = RGBColor(0xFF, 0xD4, 0x26)
    C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
    C_DGREY  = RGBColor(0x1A, 0x1A, 0x1A)
    C_MGREY  = RGBColor(0x3D, 0x3D, 0x3D)
    C_LGREY  = RGBColor(0x9A, 0x9A, 0x9A)
    C_GREEN  = RGBColor(0x00, 0xC8, 0x96)
    C_RED    = RGBColor(0xFF, 0x44, 0x44)
    C_ORANGE = RGBColor(0xFF, 0x98, 0x00)
    C_BLUE   = RGBColor(0x4A, 0x9E, 0xFF)

    SW = Inches(10)   # slide width
    SH = Inches(5.625)  # slide height

    prs = Presentation()
    prs.slide_width  = SW
    prs.slide_height = SH

    blank_layout = prs.slide_layouts[6]  # completely blank

    def rgb_hex(r):
        return RGBColor((r>>16)&0xFF, (r>>8)&0xFF, r&0xFF)

    def add_rect(slide, x, y, w, h, fill_color, line_color=None):
        shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        if line_color:
            shape.line.color.rgb = line_color
            shape.line.width = Pt(0.5)
        else:
            shape.line.fill.background()
        return shape

    def add_text_box(slide, text, x, y, w, h, size=12, bold=False, color=None,
                     align=PP_ALIGN.LEFT, font_name="Arial", italic=False):
        txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf  = txb.text_frame
        tf.word_wrap = False
        p   = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = str(text)[:120]
        run.font.name = font_name
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        if color:
            run.font.color.rgb = color
        return txb

    def brand_bar(slide):
        """Left yellow accent bar + BL mark."""
        add_rect(slide, 0, 0, 0.07, 5.625, C_YELLOW)
        box = add_rect(slide, 0.18, 0.22, 0.6, 0.36, C_YELLOW)
        add_text_box(slide, "BL", 0.18, 0.22, 0.6, 0.36, size=9, bold=True,
                     color=C_BLACK, align=PP_ALIGN.CENTER)

    def slide_title(slide, title, y=0.2):
        add_text_box(slide, title, 0.25, y, 9.5, 0.5, size=16, bold=True, color=C_WHITE)
        add_rect(slide, 0.25, y+0.55, 9.5, 0.02, C_MGREY)

    def section_divider(title, subtitle=""):
        s = prs.slides.add_slide(blank_layout)
        s.background.fill.solid()
        s.background.fill.fore_color.rgb = C_BLACK
        add_rect(s, 0, 0, 0.07, 5.625, C_YELLOW)
        add_text_box(s, "BL", 0.18, 0.22, 0.6, 0.36, size=9, bold=True, color=C_BLACK, align=PP_ALIGN.CENTER)
        add_text_box(s, title,    1.1, 1.8, 8.0, 1.0, size=34, bold=True, color=C_WHITE)
        if subtitle:
            add_text_box(s, subtitle, 1.1, 2.9, 8.0, 0.5, size=13, color=C_LGREY)
        return s

    def content_slide(title):
        s = prs.slides.add_slide(blank_layout)
        s.background.fill.solid()
        s.background.fill.fore_color.rgb = C_DGREY
        brand_bar(s)
        slide_title(s, title)
        return s

    def add_check_result_table(slide, entries, y_start=0.9, max_rows=14):
        rows = entries[:max_rows]
        ncols = 2
        col_w = [5.0, 4.5]
        # header
        for ci, (label, w) in enumerate(zip(["Check", "Result"], col_w)):
            x = 0.25 + sum(col_w[:ci])
            add_rect(slide, x, y_start, w, 0.28, C_BLACK)
            add_text_box(slide, label, x+0.05, y_start+0.03, w-0.1, 0.22,
                         size=8, bold=True, color=C_YELLOW)
        # rows
        for ri, e in enumerate(rows):
            y = y_start + 0.28 + ri * 0.27
            bg = C_DGREY if ri % 2 == 0 else RGBColor(0x22, 0x22, 0x22)
            check  = str(e.get("Check", ""))[:60]
            result = e.get("Result", "")
            if isinstance(result, (list, dict)):
                result = str(result)[:80]
            result = str(result)[:80]
            rc = C_GREEN if "✅" in result else C_RED if "❌" in result else C_LGREY
            for ci, (val, w, vc) in enumerate(zip([check, result], col_w, [C_WHITE, rc])):
                x = 0.25 + sum(col_w[:ci])
                add_rect(slide, x, y, w, 0.25, bg)
                add_text_box(slide, val, x+0.05, y+0.03, w-0.1, 0.20, size=8, color=vc)

    def add_simple_table(slide, headers, rows, y_start=0.9, max_rows=14):
        rows = rows[:max_rows]
        ncols = len(headers)
        col_w = [9.5 / ncols] * ncols
        # header
        for ci, (h, w) in enumerate(zip(headers, col_w)):
            x = 0.25 + ci * w
            add_rect(slide, x, y_start, w, 0.28, C_BLACK)
            add_text_box(slide, h, x+0.05, y_start+0.03, w-0.1, 0.22, size=8, bold=True, color=C_YELLOW)
        # rows
        for ri, row in enumerate(rows):
            y = y_start + 0.28 + ri * 0.27
            bg = C_DGREY if ri % 2 == 0 else RGBColor(0x22, 0x22, 0x22)
            for ci, (val, w) in enumerate(zip(row, col_w)):
                x = 0.25 + ci * w
                add_rect(slide, x, y, w, 0.25, bg)
                add_text_box(slide, str(val)[:60], x+0.05, y+0.03, w-0.1, 0.20, size=8, color=C_WHITE)

    # ── Build summary data ─────────────────────────────────────────────────
    summary_items = []
    if audit_data.get("Duplicate Transactions"):
        summary_items.append(("❌ Duplicate transactions detected", C_RED))
    if audit_data.get("Transaction Where Item Data Missing"):
        summary_items.append(("❌ Purchase events with missing item names", C_RED))
    if any(str(e.get("Result","")).startswith("❌") for e in audit_data.get("PII Check",[])):
        summary_items.append(("❌ Potential PII detected in page paths", C_RED))
    lp_pct_raw = next((e.get("Result","0%") for e in audit_data.get("Landing Page Analysis",[]) if e.get("Check")=="Landing Page (not set) %"), "0%")
    ua_pct_raw = next((e.get("Result","0%") for e in audit_data.get("Channel Grouping Analysis",[]) if e.get("Check")=="Unassigned %"), "0%")
    try:
        if float(str(lp_pct_raw).replace("%","")) > 10:
            summary_items.append((f"⚠️ Landing Page (not set) rate: {lp_pct_raw}", C_ORANGE))
    except: pass
    try:
        if float(str(ua_pct_raw).replace("%","")) > 10:
            summary_items.append((f"⚠️ Unassigned traffic: {ua_pct_raw}", C_ORANGE))
    except: pass
    if not summary_items:
        summary_items.append(("✅ GA4 property appears healthy across all checks", C_GREEN))

    def g(sec, chk):
        return next((str(e.get("Result","—")) for e in audit_data.get(sec,[]) if e.get("Check")==chk), "—")

    kpis = [
        ("Time Zone",      g("Property Details","Time Zone"),      C_BLUE),
        ("Currency",       g("Property Details","Currency"),       C_YELLOW),
        ("Retention",      g("Property Details","Retention Period"),C_BLUE),
        ("Key Events",     g("GA4 Property Limits","Key Events Used"), C_GREEN),
        ("LP (not set) %", str(lp_pct_raw), C_RED if float(str(lp_pct_raw).replace("%","") or 0)>10 else C_GREEN),
        ("Unassigned %",   str(ua_pct_raw), C_RED if float(str(ua_pct_raw).replace("%","") or 0)>10 else C_GREEN),
    ]

    # ── SLIDE 1: Cover ─────────────────────────────────────────────────────
    cover = prs.slides.add_slide(blank_layout)
    cover.background.fill.solid()
    cover.background.fill.fore_color.rgb = C_BLACK
    add_rect(cover, 0, 0, 0.07, 5.625, C_YELLOW)
    add_rect(cover, 0.07, 3.8, 9.93, 0.06, C_YELLOW)
    add_text_box(cover, "BL", 0.18, 0.22, 0.6, 0.36, size=9, bold=True, color=C_BLACK, align=PP_ALIGN.CENTER)
    add_text_box(cover, "GA4 AUDIT REPORT", 0.35, 1.2, 9.0, 0.55, size=10, bold=True, color=C_YELLOW)
    add_text_box(cover, prop_name[:60],      0.35, 1.9, 9.0, 1.0,  size=36, bold=True, color=C_WHITE)
    add_text_box(cover, date_range,          0.35, 3.1, 6.0, 0.5,  size=13, color=C_LGREY)

    # ── SLIDE 2: Audit Summary ─────────────────────────────────────────────
    sum_slide = content_slide("Audit Summary")
    for i, (text, col) in enumerate(summary_items[:7]):
        y = 1.0 + i * 0.62
        add_rect(sum_slide, 0.25, y, 9.5, 0.52, RGBColor(0x11,0x11,0x11))
        add_rect(sum_slide, 0.25, y, 0.05, 0.52, col)
        add_text_box(sum_slide, text[:100], 0.45, y+0.1, 9.0, 0.35, size=11, color=C_WHITE)

    # ── SLIDE 3: KPI Dashboard ─────────────────────────────────────────────
    kpi_slide = content_slide("Property Dashboard")
    positions = [(0.25,1.0),(3.55,1.0),(6.85,1.0),(0.25,2.7),(3.55,2.7),(6.85,2.7)]
    for i, (label, value, col) in enumerate(kpis):
        if i >= len(positions): break
        px, py = positions[i]
        add_rect(kpi_slide, px, py, 3.0, 1.3, RGBColor(0x11,0x11,0x11))
        add_rect(kpi_slide, px, py, 0.05, 1.3, col)
        add_text_box(kpi_slide, label.upper(), px+0.12, py+0.08, 2.75, 0.22, size=7, bold=True, color=C_LGREY)
        add_text_box(kpi_slide, str(value)[:25], px+0.12, py+0.38, 2.75, 0.6, size=19, bold=True, color=col)

    # ── SLIDE 4: Property Details ──────────────────────────────────────────
    section_divider("Property Details", "Settings · Streams · Limits")
    prop_s = content_slide("Property Details & Streams")
    add_check_result_table(prop_s, audit_data.get("Property Details",[]), y_start=0.9, max_rows=10)

    streams = audit_data.get("Streams Configuration",[])
    if streams:
        st_s = content_slide("Streams Configuration")
        add_check_result_table(st_s, streams, y_start=0.9, max_rows=10)

    limits = audit_data.get("GA4 Property Limits",[])
    if limits:
        lim_s = content_slide("GA4 Property Limits")
        add_check_result_table(lim_s, limits, y_start=0.9, max_rows=10)

    # ── SLIDE 5: Custom Dimensions ─────────────────────────────────────────
    section_divider("Custom Dimensions", "Event · User · Item scoped")
    for scope in ["Event Scoped","User Scoped","Item Scoped"]:
        dims = audit_data.get(f"Custom Dimensions - {scope}",[])
        if dims:
            ds = content_slide(f"Custom Dimensions — {scope}")
            rows = [(e.get("Check",""), e.get("Result",{}).get("Parameter Name",""), scope.split()[0]) for e in dims]
            add_simple_table(ds, ["Display Name","Parameter Name","Scope"], rows, y_start=0.9, max_rows=14)

    # ── SLIDE 6: Event Inventory ───────────────────────────────────────────
    section_divider("Events & Key Events", "What's being tracked")
    events = audit_data.get("GA4 Events",[])
    if events:
        ev_s = content_slide("Event Inventory (Top 20)")
        rows = [(e.get("Check",""), e.get("Result","")) for e in events[:20]]
        add_simple_table(ev_s, ["Event Name","Event Count"], rows, y_start=0.9, max_rows=14)

    key_events = audit_data.get("Key Event Details",[])
    if key_events:
        ke_s = content_slide("Key Events (Conversions)")
        rows = [(e.get("Check",""), e.get("Result",{}).get("Create Time",""), e.get("Result",{}).get("Counting Method","")) for e in key_events[:12]]
        add_simple_table(ke_s, ["Event Name","Created","Counting Method"], rows, y_start=0.9, max_rows=12)

    # ── SLIDE 7: PII Check ─────────────────────────────────────────────────
    pii = audit_data.get("PII Check",[])
    if pii:
        section_divider("PII & Data Quality", "")
        pii_s = content_slide("PII Check")
        add_check_result_table(pii_s, pii, y_start=0.9, max_rows=12)

    # ── SLIDE 8: Transactions ──────────────────────────────────────────────
    txns = audit_data.get("Transactions",[])
    if txns:
        section_divider("E-commerce & Transactions", "")
        tx_s = content_slide("Transaction Health")
        add_check_result_table(tx_s, txns, y_start=0.9, max_rows=12)

    dup = audit_data.get("Duplicate Transactions",[])
    if dup:
        dup_s = content_slide("Duplicate Transactions")
        rows = [(e.get("transactionId",""), e.get("count","")) for e in dup[:14]]
        add_simple_table(dup_s, ["Transaction ID","Count"], rows, y_start=0.9, max_rows=14)

    # ── SLIDE 9: Traffic Quality ───────────────────────────────────────────
    section_divider("Traffic Quality", "Landing pages · Channel grouping")
    lp_anal = audit_data.get("Landing Page Analysis",[])
    if lp_anal:
        lpa_s = content_slide("Landing Page Analysis")
        add_check_result_table(lpa_s, lp_anal, y_start=0.9, max_rows=8)

    lp_data = audit_data.get("Landing Page Data",[])
    if lp_data:
        lpd_s = content_slide("Landing Page Data (Top 14)")
        rows = [(e.get("Landing Page","")[:50], e.get("Sessions","")) for e in lp_data[:14]]
        add_simple_table(lpd_s, ["Landing Page","Sessions"], rows, y_start=0.9, max_rows=14)

    ch_anal = audit_data.get("Channel Grouping Analysis",[])
    if ch_anal:
        cha_s = content_slide("Channel Grouping Analysis")
        add_check_result_table(cha_s, ch_anal, y_start=0.9, max_rows=8)

    ch_data = audit_data.get("Channel Grouping Data",[])
    if ch_data:
        chd_s = content_slide("Channel Grouping Data")
        rows = [(e.get("Channel Group",""), e.get("Sessions","")) for e in ch_data[:14]]
        add_simple_table(chd_s, ["Channel Group","Sessions"], rows, y_start=0.9, max_rows=14)

    ua_data = audit_data.get("Unassigned Source/Medium Data",[])
    if ua_data:
        ua_s = content_slide("Unassigned Source/Medium")
        rows = [(e.get("Source",""), e.get("Medium",""), e.get("Sessions","")) for e in ua_data[:14]]
        add_simple_table(ua_s, ["Source","Medium","Sessions"], rows, y_start=0.9, max_rows=14)

    # ── END CARD ───────────────────────────────────────────────────────────
    end = prs.slides.add_slide(blank_layout)
    end.background.fill.solid()
    end.background.fill.fore_color.rgb = C_BLACK
    add_rect(end, 0, 0,      SW.inches, 0.07, C_YELLOW)
    add_rect(end, 0, 5.555,  SW.inches, 0.07, C_YELLOW)
    add_text_box(end, "BL",            4.65, 1.7,  0.7, 0.38, size=9, bold=True, color=C_BLACK, align=PP_ALIGN.CENTER)
    add_text_box(end, "Brainlabs",     3.5,  2.2,  3.0, 0.55, size=22, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text_box(end, "GA4 Audit Tool",3.2,  2.82, 3.6, 0.4,  size=12, color=C_LGREY, align=PP_ALIGN.CENTER)

    # ── Stream to response ─────────────────────────────────────────────────
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    filename = f"GA4_Audit_{prop_name.replace(' ','_')}.pptx"
    return StreamingResponse(
        buf,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": f'attachment; filename="{filename}"'},
    )


# ── Claude AI helper for PPTX Obs/Rec/Impact ─────────────────────────────
import httpx, json as _json

def _claude_analyse(section_name: str, section_data: dict, property_context: dict) -> dict:
    """
    Calls Claude claude-sonnet-4-20250514 with real audit data for a specific section.
    Returns {"observation": str, "recommendation": str, "impact": str, "badge": str}
    badge is one of: Pass | Fail | High | Medium | Need to be discussed
    """
    prompt = f"""You are a GA4 analytics expert at Brainlabs, a performance marketing agency.
You are writing one slide of a GA4 Audit Report for a client.

Property context:
{_json.dumps(property_context, indent=2)}

Audit section: {section_name}
Real data from this section:
{_json.dumps(section_data, indent=2)}

Write a professional GA4 audit finding for this section. Be specific — reference the ACTUAL values from the data above, not generic advice. If a value is missing or unclear, say so explicitly.

Respond ONLY with a JSON object, no markdown, no explanation:
{{
  "observation": "2-3 sentences. State exactly what the data shows. Reference specific values.",
  "recommendation": "2-3 sentences. Specific actionable steps. Reference what needs to change.",
  "impact": "2 sentences. Business impact of fixing this. Be concrete.",
  "badge": "one of: Pass | Fail | High | Need to be discussed"
}}

Badge rules:
- Pass: everything is correctly configured
- Fail: critical misconfiguration that must be fixed
- High: same as Fail (use Fail)
- Need to be discussed: data not visible, needs client input, or ambiguous
"""
    try:
        resp = httpx.post(
            "https://api.anthropic.com/v1/messages",
            headers={
                "x-api-key": os.environ.get("ANTHROPIC_API_KEY", ""),
                "anthropic-version": "2023-06-01",
                "content-type": "application/json",
            },
            json={
                "model": "claude-sonnet-4-20250514",
                "max_tokens": 600,
                "messages": [{"role": "user", "content": prompt}],
            },
            timeout=30.0,
        )
        text = resp.json()["content"][0]["text"].strip()
        # Strip markdown fences if present
        if text.startswith("```"):
            text = "\n".join(text.split("\n")[1:])
        if text.endswith("```"):
            text = "\n".join(text.split("\n")[:-1])
        return _json.loads(text)
    except Exception as e:
        # Graceful fallback — section-specific static text so PPTX always builds
        static_fallbacks = {
            "Data Retention": {
                "observation": f"Data retention is set to: {section_data.get('current_retention', 'unknown')}. The maximum for Standard properties is 14 months.",
                "recommendation": "Navigate to Admin → Data Settings → Data Retention and set both Event data and User data to 14 months. Enable 'Reset on new user activity'.",
                "impact": "14-month retention restores year-on-year comparison capability and supports audience seed sizes for remarketing.",
                "badge": "Pass" if "14" in str(section_data.get("current_retention","")) else "Fail",
            },
            "Consent Management": {
                "observation": "Consent Management configuration is not visible in the GA4 dashboard export. Must be verified directly in GTM.",
                "recommendation": "Confirm a CMP (e.g. Cookiebot, OneTrust) is active. Verify Consent Mode v2 default and update calls fire through GTM for all tags.",
                "impact": "Without Consent Mode v2, behavioural modelling for unconsented users is unavailable and the property risks GDPR/DPDP non-compliance.",
                "badge": "Need to be discussed",
            },
            "Google Signals": {
                "observation": "Google Signals status is not directly visible in the GA4 dashboard export. User Data Collection acknowledgement state was checked.",
                "recommendation": "Confirm Google Signals is enabled on each data stream under Admin → Data Collection. Required for cross-device user counts and Ads personalisation.",
                "impact": "Google Signals enables cross-device reporting and audience export to Google Ads. Disabling it reduces remarketing reach.",
                "badge": "Need to be discussed",
            },
            "PII Check": {
                "observation": f"{section_data.get('issues_found', 0)} PII issue(s) detected in page paths and event parameters. UTMs, click IDs, and GA4-redacted values were excluded.",
                "recommendation": "Remove any personally identifiable information from URL parameters and event data. Enable GA4 Data Redaction as an additional safeguard.",
                "impact": "Collecting PII in GA4 violates Google's Terms of Service and exposes the business to GDPR/DPDP compliance risk.",
                "badge": "Pass" if section_data.get("issues_found", 0) == 0 else "Fail",
            },
            "Custom Dimensions": {
                "observation": f"{section_data.get('event_scoped_count',0)} event-scoped, {section_data.get('user_scoped_count',0)} user-scoped, and {section_data.get('item_scoped_count',0)} item-scoped custom dimensions are configured.",
                "recommendation": "Introduce user-level dimensions (login_status, user_type) and event-level dimensions for content categorisation. Ensure parameter names match your dataLayer exactly.",
                "impact": "Custom dimensions enable segmentation beyond GA4 defaults and are critical for understanding user behaviour specific to your business model.",
                "badge": "Pass" if (section_data.get('event_scoped_count',0) + section_data.get('user_scoped_count',0)) > 0 else "Fail",
            },
            "Key Events (Conversions)": {
                "observation": f"{section_data.get('total_key_events', 0)} key event(s) configured. Key events drive Smart Bidding and conversion reporting.",
                "recommendation": "Review all business-critical actions and mark as key events: form_submit, newsletter_signup, add_to_cart, begin_checkout, purchase.",
                "impact": "Key events drive conversion reporting and Smart Bidding signals in Google Ads. Sparse configuration limits marketing optimisation.",
                "badge": "Pass" if section_data.get("total_key_events", 0) >= 3 else "Fail",
            },
            "Transaction Health & Data Duplication": {
                "observation": f"{section_data.get('duplicate_count', 0)} duplicate transaction IDs detected. Duplicates inflate purchase counts and revenue metrics.",
                "recommendation": "Ensure purchase events fire only once per order. Verify transaction_id is unique per purchase and not re-triggered on page refresh. Test via GA4 DebugView.",
                "impact": "Duplicate transactions directly inflate ROAS figures and purchase counts, leading to inaccurate eCommerce reporting and misallocated budget.",
                "badge": "Pass" if section_data.get("duplicate_count", 0) == 0 else "Fail",
            },
            "Event Inventory & Tracking": {
                "observation": f"{section_data.get('total_events', 0)} events tracked in the audit date range. Event names should follow snake_case GA4 naming conventions.",
                "recommendation": "Ensure all key user interactions are tracked as events. Add micro-conversions such as form_submit, cta_click, and scroll_depth where missing.",
                "impact": "Comprehensive event tracking is the foundation of GA4 reporting. Missing events create blind spots in funnel analysis and audience building.",
                "badge": "Pass" if section_data.get("total_events", 0) > 5 else "Need to be discussed",
            },
            "Landing Page Analysis": {
                "observation": f"Landing Page (not set) rate is {section_data.get('not_set_percent','0%')} of total sessions. A high rate indicates session_start may not be firing on all entry points.",
                "recommendation": "Investigate pages where session_start is not firing. Common causes: SPAs without history-change tracking, iframes, and redirect chains that drop the referrer.",
                "impact": "A high (not set) rate degrades acquisition and content reporting, making it impossible to measure where users first enter the site.",
                "badge": "Pass" if float(str(section_data.get("not_set_percent","0%")).replace("%","") or 0) <= 10 else "Fail",
            },
            "Unassigned Traffic & Channel Grouping": {
                "observation": f"Unassigned traffic is {section_data.get('unassigned_percent','0%')} of total sessions. Unassigned sessions cannot be attributed to a marketing channel.",
                "recommendation": "Review UTM parameter consistency across all campaigns. Configure Custom Channel Groups in GA4. Validate Google Ads auto-tagging is enabled.",
                "impact": "Unassigned traffic creates attribution gaps, making it impossible to accurately measure which channels drive value and affecting budget allocation.",
                "badge": "Pass" if float(str(section_data.get("unassigned_percent","0%")).replace("%","") or 0) <= 10 else "Need to be discussed",
            },
            "Product Linking (Google Ads & Firebase)": {
                "observation": f"Google Ads: {section_data.get('google_ads','Not checked')}. Firebase: {section_data.get('firebase','Not checked')}.",
                "recommendation": "Link Google Ads for Smart Bidding signals and audience sharing. Link Firebase for app + web cross-platform reporting. Also consider BigQuery, Search Console, and DV360.",
                "impact": "Product links unlock data synergy between GA4 and the Google Marketing Platform. Unlinked products limit advanced analysis and campaign optimisation.",
                "badge": "Pass" if "✅" in str(section_data.get("google_ads","")) else "Need to be discussed",
            },
        }
        fallback = static_fallbacks.get(section_name, {
            "observation": f"Section: {section_name}. Data reviewed: {str(section_data)[:150]}",
            "recommendation": "Review this section against GA4 best practices and your business requirements.",
            "impact": "Proper GA4 configuration ensures accurate data collection, better reporting, and improved marketing decisions.",
            "badge": "Need to be discussed",
        })
        return fallback


@app.post("/export/pptx-bl")
def export_pptx_bl(request: Request, body: dict = None):
    """
    Pixel-perfect replica of the BL GA4 Audit template.
    Uses Claude AI to generate real Observation / Recommendation / Impact for every finding.
    """
    import io
    from datetime import date
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.oxml.ns import qn
    from lxml import etree
    from fastapi.responses import StreamingResponse

    get_user_credentials(request)
    if not body:
        return {"success": False, "error": "Body required."}

    audit_data = body.get("data", {})
    prop_name  = body.get("property_name", "GA4 Property")
    date_range = body.get("date_range", "")
    today_str  = date.today().strftime("%B %Y")

    # ── Property context passed to every Claude call ───────────────────────
    prop_ctx = {
        "property_name":  prop_name,
        "date_range":     date_range,
        "property_details": audit_data.get("Property Details", []),
        "service_level":  next((e.get("Result","") for e in audit_data.get("Property Details",[]) if e.get("Check")=="Service Level"), "Unknown"),
    }

    # ── EXACT COLOURS FROM TEMPLATE ────────────────────────────────────────
    C_BG        = RGBColor(0xFF,0xFE,0xF7)
    C_BLACK     = RGBColor(0x00,0x00,0x00)
    C_WHITE     = RGBColor(0xFF,0xFF,0xFF)
    C_YELLOW_HL = RGBColor(0xFF,0xE7,0x70)
    C_BLUE_LIGHT= RGBColor(0xEB,0xF9,0xFF)
    C_BLUE_SEC  = RGBColor(0x80,0xDB,0xFF)
    C_YELLOW_NT = RGBColor(0xFF,0xF5,0xC2)
    C_GREEN     = RGBColor(0x4C,0xAF,0x50)
    C_AMBER     = RGBColor(0xFF,0xC1,0x07)
    C_RED       = RGBColor(0xE5,0x39,0x35)
    C_LGREY     = RGBColor(0x99,0x99,0x99)

    prs = Presentation()
    prs.slide_width  = Emu(12191695)
    prs.slide_height = Emu(6858000)
    blank = prs.slide_layouts[6]

    # ── SHAPE / TEXT HELPERS ───────────────────────────────────────────────
    def _i(v): return Inches(v)

    def add_rect(slide, x, y, w, h, fill, line=None, line_pt=0.75, shape='rect'):
        sp = slide.shapes.add_shape(1, _i(x), _i(y), _i(w), _i(h))
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
        if line: sp.line.color.rgb = line; sp.line.width = Pt(line_pt)
        else: sp.line.fill.background()
        if shape in ('ellipse','roundRect'):
            spPr = sp._element.find(qn('p:spPr'))
            old  = spPr.find(qn('a:prstGeom'))
            if old is not None: spPr.remove(old)
            ng = etree.SubElement(spPr, qn('a:prstGeom')); ng.set('prst', shape)
            avLst = etree.SubElement(ng, qn('a:avLst'))
            if shape == 'roundRect':
                gd = etree.SubElement(avLst, qn('a:gd'))
                gd.set('name','adj'); gd.set('fmla','val 16667')
        return sp

    def add_tb(slide, text, x, y, w, h, size=11, bold=False, italic=False,
               color=None, align=PP_ALIGN.LEFT, wrap=True):
        txb = slide.shapes.add_textbox(_i(x),_i(y),_i(w),_i(h))
        tf  = txb.text_frame; tf.word_wrap = wrap
        p   = tf.paragraphs[0]; p.alignment = align
        run = p.add_run(); run.text = str(text)
        run.font.name = "Calibri"; run.font.size = Pt(size)
        run.font.bold = bold; run.font.italic = italic
        run.font.color.rgb = color or C_BLACK
        return txb

    def add_badge(slide, text, x, y, w, h, bg, tc=C_WHITE, size=12):
        sp = add_rect(slide, x, y, w, h, bg, shape='roundRect')
        sp.line.fill.background()
        tf = sp.text_frame; tf.word_wrap = True
        p  = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        run = p.add_run(); run.text = text
        run.font.name="Calibri"; run.font.size=Pt(size)
        run.font.bold=True; run.font.color.rgb=tc
        return sp

    def add_conn(slide, x1, y1, x2, y2):
        c = slide.shapes.add_connector(1,_i(x1),_i(y1),_i(x2),_i(y2))
        c.line.color.rgb = C_BLACK; c.line.width = Pt(0.75)

    def set_bg(slide, color):
        bg=slide.background; f=bg.fill; f.solid(); f.fore_color.rgb=color

    def bottom_line(slide):
        add_rect(slide, 0.30, 7.15, 12.73, 0.03, C_LGREY)

    def page_num(slide, n):
        add_tb(slide, str(n), 0.30, 7.20, 0.50, 0.25, size=10)

    def badge_settings(badge_text):
        """Return (bg_color, text_color) for a badge label."""
        t = badge_text.lower()
        if "pass"  in t: return C_GREEN, C_WHITE
        if "fail"  in t or "high" in t: return C_RED, C_WHITE
        return C_AMBER, C_BLACK   # TBD / Need to be discussed / Medium

    # ── DATA HELPERS ───────────────────────────────────────────────────────
    def gv(section, check, default="[To be confirmed]"):
        return next((str(e.get("Result","")) for e in audit_data.get(section,[]) if e.get("Check")==check), default)

    def pct(v):
        try: return float(str(v).replace("%","").strip())
        except: return 0.0

    # Pre-compute key metrics
    ret_val    = gv("Property Details","Retention Period","Two Months")
    lp_pct_raw = gv("Landing Page Analysis","Landing Page (not set) %","0%")
    ua_pct_raw = gv("Channel Grouping Analysis","Unassigned %","0%")
    dup_raw    = gv("Transactions","Duplicate Transaction Count","0")
    try: dup_count = int(str(dup_raw).replace(",",""))
    except: dup_count = 0
    pii_issues  = [e for e in audit_data.get("PII Check",[]) if "❌" in str(e.get("Result",""))]
    event_dims  = audit_data.get("Custom Dimensions - Event Scoped",[])
    user_dims   = audit_data.get("Custom Dimensions - User Scoped",[])
    item_dims   = audit_data.get("Custom Dimensions - Item Scoped",[])
    total_dims  = len(event_dims)+len(user_dims)+len(item_dims)
    ke_list     = audit_data.get("Key Event Details",[])
    events      = audit_data.get("GA4 Events",[])

    # ── AI GENERATION — call Claude for each section ───────────────────────
    # Each call sends the real data; results used directly in the PPTX slides
    ai = {}

    ai["retention"] = _claude_analyse("Data Retention", {
        "current_retention": ret_val,
        "service_level": prop_ctx["service_level"],
        "all_checks": audit_data.get("Property Details",[]),
    }, prop_ctx)

    ai["consent"] = _claude_analyse("Consent Management", {
        "note": "Consent Mode configuration is not visible in the GA4 dashboard export. Must be verified in GTM.",
        "property_details": audit_data.get("Property Details",[]),
    }, prop_ctx)

    ai["signals"] = _claude_analyse("Google Signals", {
        "property_details": audit_data.get("Property Details",[]),
        "note": "Google Signals toggle visibility depends on stream configuration",
    }, prop_ctx)

    ai["pii"] = _claude_analyse("PII Check", {
        "pii_checks": audit_data.get("PII Check",[]),
        "issues_found": len(pii_issues),
        "issue_details": [e.get("Result","") for e in pii_issues[:5]],
    }, prop_ctx)

    ai["custom_dims"] = _claude_analyse("Custom Dimensions", {
        "event_scoped_count": len(event_dims),
        "user_scoped_count":  len(user_dims),
        "item_scoped_count":  len(item_dims),
        "event_dims": [{"name":e.get("Check",""),"param":e.get("Result",{}).get("Parameter Name","")} for e in event_dims[:10]],
        "user_dims":  [{"name":e.get("Check",""),"param":e.get("Result",{}).get("Parameter Name","")} for e in user_dims[:5]],
    }, prop_ctx)

    ai["key_events"] = _claude_analyse("Key Events (Conversions)", {
        "total_key_events": len(ke_list),
        "key_events": [{"name":e.get("Check",""), "counting_method":e.get("Result",{}).get("Counting Method","")} for e in ke_list[:10]],
    }, prop_ctx)

    ai["transactions"] = _claude_analyse("Transaction Health & Data Duplication", {
        "duplicate_count": dup_count,
        "transaction_checks": audit_data.get("Transactions",[]),
        "duplicate_examples": audit_data.get("Duplicate Transactions",[])[:5],
    }, prop_ctx)

    ai["events"] = _claude_analyse("Event Inventory & Tracking", {
        "total_events": len(events),
        "top_events": [{"name":e.get("Check",""), "count":e.get("Result","")} for e in events[:15]],
    }, prop_ctx)

    ai["landing_page"] = _claude_analyse("Landing Page Analysis", {
        "not_set_percent": lp_pct_raw,
        "total_sessions": gv("Landing Page Analysis","Total Sessions","unknown"),
        "top_landing_pages": audit_data.get("Landing Page Data",[])[:8],
    }, prop_ctx)

    ai["unassigned"] = _claude_analyse("Unassigned Traffic & Channel Grouping", {
        "unassigned_percent": ua_pct_raw,
        "unassigned_sources": audit_data.get("Unassigned Source/Medium Data",[])[:8],
        "channel_data": audit_data.get("Channel Grouping Data",[])[:8],
    }, prop_ctx)

    ai["product_links"] = _claude_analyse("Product Linking (Google Ads & Firebase)", {
        "google_ads": gv("Product Links","Google Ads","Not checked"),
        "firebase":   gv("Product Links","Firebase","Not checked"),
        "all_product_links": audit_data.get("Product Links",[]),
    }, prop_ctx)

    # ── AUDIT OVERVIEW TABLE — also AI-driven badge + comment ─────────────
    def _ov_badge(val_str):
        v = str(val_str)
        if "✅" in v: return (C_GREEN, C_WHITE)
        if "❌" in v: return (C_RED,   C_WHITE)
        return (C_AMBER, C_BLACK)

    def _ov_entry(section_key, check_key, fallback="Not checked"):
        entry = next((e for e in audit_data.get(section_key,[]) if e.get("Check")==check_key), None)
        return str(entry.get("Result", fallback)) if entry else fallback

    # Overview rows now use REAL data from the new APIs
    overview_rows_1 = [
        ("Consent Management",    ai["consent"]["observation"][:120],     badge_settings(ai["consent"]["badge"])),
        ("Google Signals",        _ov_entry("Google Signals","State"),    _ov_badge(_ov_entry("Google Signals","State"))),
        ("Data Retention",        ai["retention"]["observation"][:120],   badge_settings(ai["retention"]["badge"])),
        ("Data Filters / Internal traffic", "Filter state not in dashboard. Confirm internal traffic exclusion is Active, not Testing.", (C_AMBER, C_BLACK)),
        ("Cross-domain measurement", "Not visible in dashboard. Confirm if cross-domain journeys exist and configure if required.", (C_AMBER, C_BLACK)),
        ("Session timeout",       "Not visible in dashboard. Confirm defaults (30 min / 10 sec) are in place.", (C_AMBER, C_BLACK)),
        ("Reporting Identity",    _ov_entry("Reporting Identity","Reporting Identity","⚠️ Not fetched"), _ov_badge(_ov_entry("Reporting Identity","Reporting Identity"))),
        ("Attribution Settings",  _ov_entry("Attribution Settings","Reporting Attribution Model","⚠️ Not fetched"), _ov_badge(_ov_entry("Attribution Settings","Reporting Attribution Model"))),
        ("User Provided Data",    _ov_entry("User Provided Data","Collection State","⚠️ Not fetched"), _ov_badge(_ov_entry("User Provided Data","Collection State"))),
        ("PII",                   ai["pii"]["observation"][:120],          badge_settings(ai["pii"]["badge"])),
    ]
    overview_rows_2 = [
        ("User-ID",               "Not visible. Implement if login exists for cross-device tracking.", (C_AMBER, C_BLACK)),
        ("Custom Definitions",    ai["custom_dims"]["observation"][:120],  badge_settings(ai["custom_dims"]["badge"])),
        ("Key Events",            ai["key_events"]["observation"][:120],   badge_settings(ai["key_events"]["badge"])),
        ("Transaction Health",    ai["transactions"]["observation"][:120], badge_settings(ai["transactions"]["badge"])),
        ("Landing Page (not set) %", ai["landing_page"]["observation"][:120], badge_settings(ai["landing_page"]["badge"])),
        ("Unassigned traffic",    ai["unassigned"]["observation"][:120],   badge_settings(ai["unassigned"]["badge"])),
        ("Google Ads",            _ov_entry("Product Links","Google Ads","❌ Not linked")[:120], _ov_badge(_ov_entry("Product Links","Google Ads"))),
        ("Firebase",              _ov_entry("Product Links","Firebase","❌ Not linked")[:120], _ov_badge(_ov_entry("Product Links","Firebase"))),
        ("BigQuery",              _ov_entry("Product Links","BigQuery","❌ Not linked")[:120], _ov_badge(_ov_entry("Product Links","BigQuery"))),
        ("Audiences",             f"Total: {_ov_entry('Audiences','Total Audiences','?')} | Custom: {_ov_entry('Audiences','Custom Audiences','?')}", _ov_badge(_ov_entry("Audiences","Custom Audiences","0"))),
        ("GTM Configuration",     "Container details not in dashboard. Audit container quality, tags, triggers, Consent Mode.", (C_AMBER, C_BLACK)),
    ]

    # ── SLIDE BUILDERS ─────────────────────────────────────────────────────
    def make_cover():
        s = prs.slides.add_slide(blank); set_bg(s, C_BG)
        banner = add_rect(s, 3.50, 2.20, 9.00, 1.00, C_YELLOW_HL)
        tf=banner.text_frame; p=tf.paragraphs[0]; p.alignment=PP_ALIGN.LEFT
        r=p.add_run(); r.text="  GOOGLE ANALYTICS 4"
        r.font.name="Calibri"; r.font.size=Pt(14); r.font.bold=True; r.font.color.rgb=C_BLACK
        add_tb(s,"Audit Findings &\nRecommendations",3.50,3.00,9.00,2.00,size=54,bold=True,italic=True,color=C_BLACK)
        add_tb(s,f"{prop_name} · {today_str}",3.50,5.30,9.00,0.50,size=18,color=C_BLACK)
        add_rect(s,0.80,1.80,2.00,2.00,C_WHITE,C_BLACK,0.75,'roundRect')
        add_rect(s,0.80,1.80,2.00,2.00,C_WHITE,C_BLACK,0.75,'ellipse')
        add_tb(s,prop_name[:20],0.80,2.50,2.00,0.60,size=18,bold=True,color=C_BLACK,align=PP_ALIGN.CENTER)
        add_conn(s,1.80,3.80,1.80,4.80)
        add_rect(s,0.80,4.80,2.00,2.00,C_WHITE,C_BLACK,0.75,'ellipse')
        add_tb(s,"brainlabs",0.80,5.50,2.00,0.60,size=18,bold=True,italic=True,color=C_BLACK,align=PP_ALIGN.CENTER)

    def make_property_stream():
        s = prs.slides.add_slide(blank); set_bg(s, C_BG)
        add_tb(s,"Property & Stream details",0.50,0.40,9.00,0.70,size=28,bold=True,italic=True,color=C_BLACK)
        add_badge(s,"Property details",0.50,1.30,2.00,0.35,C_BLACK,C_WHITE,12)
        prop_rows=[
            ("Property name",      gv("Property Details","Display Name")),
            ("Property ID",        gv("Property Details","Property ID")),
            ("Reporting time zone",gv("Property Details","Time Zone")),
            ("Industry Category",  gv("Property Details", "Industry Category")),
            ("Currency displayed as", gv("Property Details","Currency")),
        ]
        for i,(k,v) in enumerate(prop_rows):
            y=1.75+i*0.38
            add_rect(s,0.50,y,2.50,0.38,C_BLUE_LIGHT,C_BLACK,0.5)
            add_rect(s,3.00,y,4.50,0.38,C_WHITE,C_BLACK,0.5)
            add_tb(s,k,0.55,y+0.05,2.40,0.28,size=11)
            add_tb(s,str(v)[:60],3.05,y+0.05,4.40,0.28,size=11)
        add_badge(s,"Stream details",0.50,4.00,2.00,0.35,C_BLACK,C_WHITE,12)
        streams=audit_data.get("Streams Configuration",[])
        stream_rows=[
            ("Stream name",    gv("Streams Configuration", "Stream Name")),
            ("Stream ID",      gv("Streams Configuration", "Stream ID")),
            ("Stream URL",     gv("Streams Configuration", "Stream URL")),
            ("Measurement ID", gv("Streams Configuration", "Measurement ID")),
        ]
        for i,(k,v) in enumerate(stream_rows):
            y=4.45+i*0.38
            add_rect(s,0.50,y,2.50,0.38,C_BLUE_LIGHT,C_BLACK,0.5)
            add_rect(s,3.00,y,4.50,0.38,C_WHITE,C_BLACK,0.5)
            add_tb(s,k,0.55,y+0.05,2.40,0.28,size=11)
            add_tb(s,str(v)[:60],3.05,y+0.05,4.40,0.28,size=11)
        add_rect(s,8.50,2.00,4.30,3.50,C_BLUE_LIGHT,C_BLACK,0.75,'roundRect')
        svc=gv("Property Details","Service Level","Standard")
        add_tb(s,f"Service Level: {svc}\nRetention: {ret_val}\nUser data ack: ✓",
               8.60,3.20,4.10,1.00,size=13,color=C_BLACK,align=PP_ALIGN.CENTER)
        bottom_line(s); page_num(s,2)

    def make_overview(rows, title, pn):
        s = prs.slides.add_slide(blank); set_bg(s, C_BG)
        add_tb(s,title,0.50,0.40,9.00,0.70,size=28,bold=True,italic=True,color=C_BLACK)
        add_badge(s,"Pass",         8.60,0.55,1.30,0.30,C_GREEN,C_WHITE,12)
        add_badge(s,"TBD /\nEvaluate",10.00,0.52,1.30,0.36,C_AMBER,C_BLACK,10)
        add_badge(s,"Action\nrequired",11.40,0.52,1.40,0.36,C_RED,C_WHITE,10)
        add_rect(s,0.50,1.40,2.50,0.40,C_BLACK)
        add_rect(s,3.00,1.40,7.20,0.40,C_BLACK)
        add_rect(s,10.20,1.40,2.60,0.40,C_BLACK)
        add_tb(s,"  Audited Section",0.50,1.43,2.50,0.34,size=12,bold=True,color=C_WHITE)
        add_tb(s,"  Comments | Recommendations",3.00,1.43,7.20,0.34,size=12,bold=True,color=C_WHITE)
        add_tb(s,"  Status",10.20,1.43,2.60,0.34,size=12,bold=True,color=C_WHITE)
        for i,(section,comment,(badge_col,badge_tc)) in enumerate(rows):
            y=1.80+i*0.45
            add_rect(s,0.50, y,2.50,0.45,C_WHITE,C_BLACK,0.25)
            add_rect(s,3.00, y,7.20,0.45,C_WHITE,C_BLACK,0.25)
            add_rect(s,10.20,y,2.60,0.45,C_WHITE,C_BLACK,0.25)
            add_tb(s,section,0.55,y+0.07,2.40,0.32,size=10)
            add_tb(s,str(comment)[:130],3.05,y+0.05,7.10,0.35,size=10)
            badge_label = "Pass" if badge_col==C_GREEN else ("Fail" if badge_col==C_RED else "TBD")
            add_badge(s,badge_label,10.70,y+0.065,1.60,0.32,badge_col,badge_tc,12)
        bottom_line(s); page_num(s,pn)

    def make_divider(title, subtitle="Observations & Recommendations"):
        s = prs.slides.add_slide(blank); set_bg(s, C_BLUE_SEC)
        add_conn(s,2.00,2.80,3.00,2.80)
        add_tb(s,title,2.00,2.90,9.50,1.30,size=48,bold=True,italic=True,color=C_BLACK)
        add_tb(s,subtitle,2.00,4.20,9.50,0.40,size=14,color=C_BLACK)
        add_conn(s,2.00,5.30,11.30,5.30)

    def make_finding(title, ai_result, note_text="", pn=0):
        """Build one finding slide from a Claude AI result dict."""
        observation    = ai_result.get("observation","")
        recommendation = ai_result.get("recommendation","")
        impact         = ai_result.get("impact","")
        badge          = ai_result.get("badge","Need to be discussed")
        bg, tc = badge_settings(badge)

        s = prs.slides.add_slide(blank); set_bg(s, C_BG)
        add_tb(s,title,0.50,0.45,9.00,0.75,size=28,bold=True,color=C_BLACK)
        add_rect(s,0.50,1.25,2.00,0.06,C_BLACK)
        add_badge(s,badge,11.40,0.30,1.70,0.40,bg,tc,11)

        for (box_lbl, box_txt, by, bh) in [
            ("Observation",    observation,    1.50, 1.55),
            ("Recommendation", recommendation, 3.20, 1.55),
            ("Impact",         impact,         5.20, 1.70),
        ]:
            add_rect(s,0.50,by,6.20,bh,C_WHITE,C_BLACK,0.75,'roundRect')
            add_tb(s,box_lbl,0.65,by+0.05,5.90,0.38,size=13,bold=True)
            add_tb(s,str(box_txt)[:350],0.65,by+0.43,5.90,bh-0.52,size=10.5,wrap=True)

        add_rect(s,7.00,1.50,5.80,4.20,C_BLUE_LIGHT,C_BLACK,0.75,'roundRect')
        add_tb(s,"[Screenshot placeholder]",7.10,3.10,5.60,0.70,size=13,italic=True,
               color=C_BLACK,align=PP_ALIGN.CENTER)

        if note_text:
            add_rect(s,7.00,6.00,5.80,0.70,C_YELLOW_NT,C_BLACK,0.75,'roundRect')
            add_tb(s,str(note_text)[:100],7.10,6.05,5.60,0.60,size=10)

        bottom_line(s); page_num(s,pn)

    def make_product_links_table():
        """One slide with full product links table — all products, status + ID."""
        s = prs.slides.add_slide(blank); set_bg(s, C_BG)
        add_tb(s,"Product Links",0.50,0.45,9.00,0.75,size=28,bold=True,color=C_BLACK)
        add_rect(s,0.50,1.25,2.00,0.06,C_BLACK)

        products = [
            ("Google Ads",       "Google Ads"),
            ("Firebase",         "Firebase"),
            ("BigQuery",         "BigQuery"),
            ("Search Ads 360",   "Search Ads 360"),
            ("AdSense",          "AdSense"),
        ]
        # Table header
        add_rect(s, 0.50, 1.40, 4.50, 0.38, C_BLACK)
        add_rect(s, 5.00, 1.40, 7.80, 0.38, C_BLACK)
        add_tb(s,"  Product",       0.50, 1.43, 4.50, 0.32, size=11, bold=True, color=C_WHITE)
        add_tb(s,"  Status & Details", 5.00, 1.43, 7.80, 0.32, size=11, bold=True, color=C_WHITE)

        for i,(label, key) in enumerate(products):
            y = 1.78 + i * 0.52
            entry = next((e for e in audit_data.get("Product Links",[]) if e.get("Check")==key), None)
            val   = str(entry.get("Result","❌ Not checked")) if entry else "❌ Not checked"
            linked = "✅" in val
            bg = C_BLUE_LIGHT if linked else RGBColor(0xFF,0xEB,0xEB)
            add_rect(s, 0.50, y, 4.50, 0.50, bg, C_BLACK, 0.25)
            add_rect(s, 5.00, y, 7.80, 0.50, C_WHITE,    C_BLACK, 0.25)
            add_tb(s, label,        0.60, y+0.10, 4.30, 0.32, size=11, bold=True, color=C_BLACK)
            status_color = C_GREEN if linked else (RGBColor(0xCC,0x66,0x00) if "⚠️" in val else C_RED)
            add_tb(s, val[:100],    5.10, y+0.08, 7.60, 0.35, size=10, color=status_color, wrap=True)

        # Legend
        add_rect(s, 0.50, 4.50, 0.20, 0.20, C_BLUE_LIGHT, C_BLACK, 0.5)
        add_tb(s, "= Linked", 0.75, 4.50, 2.0, 0.20, size=9, color=C_BLACK)
        add_rect(s, 3.00, 4.50, 0.20, 0.20, RGBColor(0xFF,0xEB,0xEB), C_BLACK, 0.5)
        add_tb(s, "= Not linked", 3.25, 4.50, 2.0, 0.20, size=9, color=C_BLACK)

        bottom_line(s); page_num(s, 0)

    def make_settings_finding(title, section_key, fields_to_show, ai_section_key, note_fn=None, pn=0):
        """Generic finding slide for Settings sections (Signals, Attribution, etc.)"""
        entries = [e for e in audit_data.get(section_key,[]) if e.get("Check") != "Raw Response"]
        ai_result = _claude_analyse(title, {
            "section_data": [{"check": e.get("Check"), "result": str(e.get("Result",""))[:100]} for e in entries],
        }, prop_ctx)
        note_text = note_fn(entries) if note_fn else ""
        make_finding(title, ai_result, note_text=note_text, pn=pn)

    def make_conclusion(pn):
        s = prs.slides.add_slide(blank); set_bg(s, C_BG)
        add_tb(s,"Overall Conclusion",0.50,0.40,9.00,0.70,size=28,bold=True,italic=True)
        add_rect(s,0.50,1.25,2.00,0.06,C_BLACK)

        all_rows = overview_rows_1 + overview_rows_2
        high_items   = [(n,c) for n,c,(bc,tc) in all_rows if bc == C_RED]
        medium_items = [(n,c) for n,c,(bc,tc) in all_rows if bc == C_AMBER]
        pass_items   = [(n,c) for n,c,(bc,tc) in all_rows if bc == C_GREEN]

        add_tb(s,
            f"Audit completed for {prop_name}. "
            f"{len(high_items)} action-required item(s), "
            f"{len(medium_items)} TBD item(s), "
            f"{len(pass_items)} pass(es).",
            0.50,1.45,12.0,0.60,size=11,wrap=True)

        add_badge(s,"Action Required",0.50,2.15,3.00,0.35,C_RED,C_WHITE,11)
        for i,(sec,com) in enumerate(high_items[:6]):
            y=2.60+i*0.65
            add_rect(s,0.50,y,5.70,0.60,C_WHITE,C_BLACK,0.5,'roundRect')
            add_tb(s,f"• {sec}",0.65,y+0.03,5.40,0.22,size=10,bold=True)
            add_tb(s,str(com)[:90],0.65,y+0.28,5.40,0.28,size=9,wrap=True)

        add_badge(s,"TBD / Evaluate",6.50,2.15,3.00,0.35,C_AMBER,C_BLACK,11)
        for i,(sec,com) in enumerate(medium_items[:6]):
            y=2.60+i*0.65
            add_rect(s,6.50,y,5.70,0.60,C_WHITE,C_BLACK,0.5,'roundRect')
            add_tb(s,f"• {sec}",6.65,y+0.03,5.40,0.22,size=10,bold=True)
            add_tb(s,str(com)[:90],6.65,y+0.28,5.40,0.28,size=9,wrap=True)

        bottom_line(s); page_num(s,pn)

    def make_thank_you():
        s = prs.slides.add_slide(blank); set_bg(s, C_BG)
        add_tb(s,"Thank\nyou",0.80,2.80,8.00,2.00,size=72,bold=True,italic=True)
        add_rect(s,10.00,1.50,2.00,2.00,C_WHITE,C_BLACK,0.75,'ellipse')
        add_tb(s,prop_name[:20],10.00,2.20,2.00,0.60,size=16,bold=True,align=PP_ALIGN.CENTER)
        add_conn(s,11.00,3.50,11.00,4.50)
        add_rect(s,10.00,4.50,2.00,2.00,C_WHITE,C_BLACK,0.75,'ellipse')
        add_tb(s,"brainlabs",10.00,5.20,2.00,0.60,size=16,bold=True,italic=True,align=PP_ALIGN.CENTER)

    # ── BUILD FULL DECK ────────────────────────────────────────────────────
    make_cover()
    make_property_stream()
    make_overview(overview_rows_1, "Audit overview", 3)
    make_overview(overview_rows_2, "Audit overview (continued)", 4)

    make_divider("GA4 Configurations")
    make_finding("Data Retention",     ai["retention"],
                 note_text=f"Retention currently set to '{ret_val}'", pn=6)
    make_finding("Consent Management", ai["consent"],   pn=7)
    make_finding("Google Signals",     ai["signals"],   pn=8)
    make_finding("PII Check",          ai["pii"],
                 note_text="PII scan: UTMs and click IDs excluded." if not pii_issues else f"{len(pii_issues)} PII issue(s) found",
                 pn=9)

    # ── New: Attribution, Reporting Identity, User Provided Data ──────────
    make_divider("Property Configuration", "Attribution · Reporting Identity · User Data")

    # Attribution Settings
    attr_entries = [e for e in audit_data.get("Attribution Settings",[]) if e.get("Check")!="Raw Response"]
    attr_ai = _claude_analyse("Attribution Settings", {
        "data": [{"check": e.get("Check"), "result": str(e.get("Result",""))[:100]} for e in attr_entries]
    }, prop_ctx)
    make_finding("Attribution Settings", attr_ai,
                 note_text=f"Model: {gv('Attribution Settings','Reporting Attribution Model','Data-Driven')}",
                 pn=10)

    # Reporting Identity
    ri_entries = [e for e in audit_data.get("Reporting Identity",[]) if e.get("Check")!="Raw Response"]
    ri_ai = _claude_analyse("Reporting Identity Settings", {
        "data": [{"check": e.get("Check"), "result": str(e.get("Result",""))[:100]} for e in ri_entries]
    }, prop_ctx)
    make_finding("Reporting Identity", ri_ai,
                 note_text=gv("Reporting Identity","Reporting Identity","Not fetched"),
                 pn=11)

    # User Provided Data
    upd_entries = [e for e in audit_data.get("User Provided Data",[]) if e.get("Check")!="Raw Response"]
    upd_ai = _claude_analyse("User Provided Data Collection", {
        "data": [{"check": e.get("Check"), "result": str(e.get("Result",""))[:100]} for e in upd_entries]
    }, prop_ctx)
    make_finding("User Provided Data Collection", upd_ai,
                 note_text=gv("User Provided Data","Collection State","Not fetched"),
                 pn=12)

    # Google Signals (dedicated slide with real data)
    gs_entries = [e for e in audit_data.get("Google Signals",[]) if e.get("Check")!="Raw Response"]
    gs_ai = _claude_analyse("Google Signals", {
        "state":   gv("Google Signals","State","Not fetched"),
        "consent": gv("Google Signals","Consent","Not fetched"),
        "data": [{"check": e.get("Check"), "result": str(e.get("Result",""))[:100]} for e in gs_entries]
    }, prop_ctx)
    make_finding("Google Signals", gs_ai,
                 note_text=f"State: {gv('Google Signals','State','Not fetched')} | Consent: {gv('Google Signals','Consent','Not fetched')}",
                 pn=13)

    # ── Product Links — one table slide ───────────────────────────────────
    make_divider("Product Links", "Google Ads · Firebase · BigQuery · Search Ads 360 · AdSense")
    make_product_links_table()

    # ── Audiences ─────────────────────────────────────────────────────────
    make_divider("Audiences")
    aud_ai = _claude_analyse("Audiences", {
        "total":   gv("Audiences","Total Audiences","0"),
        "default": gv("Audiences","Default Audiences","0"),
        "custom":  gv("Audiences","Custom Audiences","0"),
        "list": [e.get("Check","").replace("Audience: ","") for e in audit_data.get("Audiences",[]) if e.get("Check","").startswith("Audience:")],
    }, prop_ctx)
    make_finding("Audiences", aud_ai,
                 note_text=f"Total: {gv('Audiences','Total Audiences','?')} | Custom: {gv('Audiences','Custom Audiences','?')}",
                 pn=14)

    make_divider("GA4 Data Quality")
    make_finding("Custom Definitions", ai["custom_dims"],
                 note_text=f"{total_dims} custom dimension(s) configured", pn=10)
    make_finding("Key Events",         ai["key_events"],
                 note_text=f"{len(ke_list)} key event(s) configured", pn=11)
    make_finding("Transaction Health", ai["transactions"],
                 note_text=f"Duplicate transactions: {dup_count}", pn=12)

    make_divider("Events & Key Events")
    make_finding("Event Inventory",    ai["events"],
                 note_text=f"{len(events)} events tracked in date range", pn=13)

    make_divider("Traffic Quality")
    make_finding("Landing Page Analysis",            ai["landing_page"],
                 note_text=f"Landing Page (not set): {lp_pct_raw}", pn=14)
    make_finding("Unassigned Traffic & Channels",    ai["unassigned"],
                 note_text=f"Unassigned sessions: {ua_pct_raw}", pn=15)

    make_conclusion(16)
    make_thank_you()

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    fn = f"GA4_Audit_{prop_name.replace(' ','_')}.pptx"
    return StreamingResponse(buf,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": f'attachment; filename="{fn}"'})


@app.get("/debug-streams")
def debug_streams(
    request: Request,
    property_id: str = Query(...),
):
    """Debug endpoint — returns raw stream data exactly as the GA4 SDK sees it."""
    from google.analytics.admin import AnalyticsAdminServiceClient
    creds = get_user_credentials(request)
    try:
        admin_client = AnalyticsAdminServiceClient(credentials=creds)
        result = []
        for stream in admin_client.list_data_streams(parent=f"properties/{property_id}"):
            entry = {
                "raw_name":         stream.name,
                "stream_id":        stream.name.split("/")[-1] if stream.name else None,
                "display_name":     stream.display_name,
                "type":             str(stream.type_),
                "has_web_data":     bool(stream.web_stream_data),
                "has_android_data": bool(stream.android_app_stream_data),
                "has_ios_data":     bool(stream.ios_app_stream_data),
            }
            if stream.web_stream_data:
                entry["measurement_id"] = stream.web_stream_data.measurement_id
                entry["default_uri"]    = stream.web_stream_data.default_uri
                entry["firebase_app_id"]= stream.web_stream_data.firebase_app_id
            result.append(entry)
        return {"success": True, "streams": result}
    except Exception as e:
        return {"success": False, "error": str(e)}